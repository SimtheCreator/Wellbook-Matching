import os
import subprocess
import sys

def install_and_import():
    try:
        import pptx
    except ImportError:
        print("python-pptx not found. Installing...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        import pptx
    return pptx

pptx = install_and_import()
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullet_slide(title, body_text):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    
    lines = body_text.split('\n')
    tf.text = lines[0]
    
    for line in lines[1:]:
        if not line.strip():
            continue
        p = tf.add_paragraph()
        p.text = line.strip()
        if line.startswith("    "):
            p.level = 2
        elif line.startswith("  "):
            p.level = 1
        else:
            p.level = 0

add_title_slide("ความสัมพันธ์ของ Vital Signs กับกิจกรรม Wellness", "ข้อมูลอ้างอิงทางการแพทย์สำหรับการออกแบบ Wellness Recommendation")

add_bullet_slide("หมวดความเครียด: Heart Rate Variability (HRV)",
"""ความหมาย: ความแปรปรวนระยะเวลาหัวใจเต้น บ่งบอกสมดุลระบบประสาท (Sympathetic vs Parasympathetic)
กิจกรรมที่แนะนำ:
  - Yoga & Breathwork: ช่วยเพิ่มความผ่อนคลาย (Parasympathetic)
  - Sauna & Cold Plunge (ทำแยกกัน): ช่วยฝึกระบบประสาทรับมือกับความเครียด
ข้อควรระวัง (Red Flag):
  - หากป่วย, อดนอน หรือ Overtraining (HRV ต่ำผิดปกติ) งดทำความเย็น/ร้อนจัด
  - ผู้ป่วยโรคหัวใจ ควรหลีกเลี่ยงอุณหภูมิสุดขั้ว""")

add_bullet_slide("หมวดความเครียด: Cortisol (ฮอร์โมนความเครียด)",
"""ความหมาย: ฮอร์โมนหลั่งจากต่อมหมวกไตเมื่อเครียด หากสูงเรื้อรังทำให้เกิดการอักเสบ ระบบเผาผลาญพัง
กิจกรรมที่แนะนำ:
  - Sound Healing & Meditation: คลื่นเสียงช่วยเปลี่ยนคลื่นสมองสู่โหมดผ่อนคลาย ลดฮอร์โมน
  - Yin / Restorative Yoga: การยืดเหยียดช่วยเคลียร์ Cortisol ออกจากกระแสเลือด
ข้อควรระวัง (Red Flag):
  - ภาวะ Burnout รุนแรง (ต่อมหมวกไตล้า) หลีกเลี่ยงกิจกรรมที่ใช้แรงกายมาก (High-intensity) หรือแช่น้ำแข็งนานเกินไป""")

add_bullet_slide("หมวดการอักเสบ: C-Reactive Protein (CRP)",
"""ความหมาย: โปรตีนสะท้อนการอักเสบในร่างกาย (Systemic Inflammation)
กิจกรรมที่แนะนำ:
  - Sauna (ประจำ): ความร้อนสร้างโปรตีนซ่อมแซมเซลล์ ลดการอักเสบระดับโมเลกุล
  - Cold Plunge: ความเย็นช่วยลดการอักเสบเฉียบพลันหลังออกกำลังกาย
ข้อควรระวัง (Red Flag):
  - หากมีไข้สูง ติดเชื้อเฉียบพลัน ห้ามทำ Sauna/Cold Plunge เด็ดขาด เพราะร่างกายต้องการพักผ่อน""")

add_bullet_slide("หมวดการอักเสบ: Resting Heart Rate (RHR)",
"""ความหมาย: อัตราการเต้นของหัวใจขณะพัก ค่าต่ำบ่งบอกหัวใจแข็งแรง ฟื้นตัวดี
กิจกรรมที่แนะนำ:
  - Yoga & Pilates: ปรับสมดุลการหายใจ ลดความดันและอัตราหัวใจเต้นระยะยาว
  - Sauna: ช่วยให้หลอดเลือดยืดหยุ่น RHR ขณะหลับจะลดลง
ข้อควรระวัง (Red Flag):
  - ผู้ที่มีอาการใจสั่น (Arrhythmia) หรือทานยา Beta-blockers ต้องระวังความร้อน/เย็นจัด""")

add_bullet_slide("หมวดระบบเผาผลาญ: ความดันโลหิต (Blood Pressure)",
"""ความหมาย: แรงดันกระแสเลือด ค่าสูงทำให้หลอดเลือดแข็งตัว เสี่ยงโรคหัวใจ
กิจกรรมที่แนะนำ:
  - Meditation, Sound Healing, Yoga: มีผลวิจัยรองรับชัดเจนว่าช่วยลดความดันโลหิตได้
ข้อควรระวัง (Red Flag สำคัญมาก!):
  - ห้ามทำ Contrast Therapy ฉับพลัน (Sauna แล้วลง Cold Plunge ทันที)
  - ความร้อนขยายหลอดเลือด หากเจอเย็นจัดหลอดเลือดจะหดตัวฉับพลัน ดันความดันพุ่งสูง อันตรายมากต่อคนเป็นโรคความดันสูง""")

add_bullet_slide("หมวดระบบเผาผลาญ: น้ำตาลในเลือด (Blood Sugar)",
"""ความหมาย: ระดับน้ำตาลในเลือด สะท้อนภาวะดื้ออินซูลินและเสี่ยงเบาหวาน
กิจกรรมที่แนะนำ:
  - Pilates & Yoga: การใช้กล้ามเนื้อมัดใหญ่ช่วยดึงน้ำตาลไปใช้ เพิ่มความไวต่ออินซูลิน โดยไม่มีแรงกระแทก
ข้อควรระวัง (Red Flag):
  - ผู้ที่ใช้ยาลดน้ำตาล/ฉีดอินซูลิน ระวังการทำ Sauna นานไป หรือทำขณะท้องว่าง
  - เหงื่อออกมากทำให้ขาดน้ำและเกิดภาวะน้ำตาลตกเฉียบพลัน (Hypoglycemia)""")

output_path = "C:\\Users\\LENOVO\\Desktop\\Codex\\Wellnista\\Wellness_Research_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
